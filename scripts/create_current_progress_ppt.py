from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "scripts" / "管道焊接标准生成多智能体系统-新版.pptx"
OUTPUT = ROOT / "scripts" / "管道焊接标准生成多智能体系统-当前进展.pptx"

NAVY = RGBColor(14, 39, 54)
TEAL = RGBColor(25, 111, 116)
GREEN = RGBColor(112, 145, 76)
PALE = RGBColor(239, 245, 243)
CARD = RGBColor(255, 255, 255)
TEXT = RGBColor(43, 57, 62)
MUTED = RGBColor(100, 114, 118)
AMBER = RGBColor(196, 126, 35)
BLUE = RGBColor(56, 105, 153)


def set_font(run, size: float, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PALE


def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, size=16, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return box


def add_title(slide, index: str, title: str, subtitle: str = "") -> None:
    add_textbox(slide, index, 1.2, 0.65, 2.2, 0.5, size=13, bold=True, color=TEAL)
    add_textbox(slide, title, 1.2, 1.18, 21.5, 0.9, size=26, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, subtitle, 1.25, 2.06, 22, 0.55, size=11.5, color=MUTED)


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, accent=TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(216, 226, 224)
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(0.08), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, title, x + 0.35, y + 0.28, w - 0.6, 0.45, size=14, bold=True, color=NAVY)
    add_textbox(slide, body, x + 0.35, y + 0.9, w - 0.6, h - 1.05, size=10.5, color=TEXT)


def add_metric(slide, x: float, y: float, number: str, label: str, color=TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(4.35), Cm(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(216, 226, 224)
    add_textbox(slide, number, x, y + 0.18, 4.35, 0.75, size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + 0.15, y + 1.08, 4.05, 0.55, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_flow(slide, items: list[tuple[str, str]], y: float) -> None:
    x = 1.2
    w = 3.95
    gap = 0.45
    for idx, (title, body) in enumerate(items):
        add_card(slide, x, y, w, 2.15, title, body, accent=[TEAL, BLUE, GREEN, AMBER, TEAL][idx % 5])
        if idx < len(items) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x + w + 0.05), Cm(y + 1.08), Cm(x + w + gap - 0.08), Cm(y + 1.08))
            conn.line.color.rgb = TEAL
            conn.line.width = Pt(2)
        x += w + gap


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size=13) -> None:
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        for run in p.runs:
            set_font(run, size, False, TEXT)


def build_ppt() -> None:
    template = Presentation(str(TEMPLATE))
    prs = Presentation()
    prs.slide_width = template.slide_width
    prs.slide_height = template.slide_height
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_textbox(slide, "Pipeline Welding", 1.25, 1.1, 8, 0.55, size=15, bold=True, color=TEAL)
    add_textbox(slide, "管道焊接标准生成\n多智能体系统", 1.2, 2.0, 15, 2.1, size=34, bold=True, color=NAVY)
    add_textbox(slide, "当前阶段工作进展汇报", 1.3, 4.4, 10, 0.7, size=18, color=TEXT)
    add_metric(slide, 1.25, 6.0, "11", "PDF标准/规范解析", TEAL)
    add_metric(slide, 6.0, 6.0, "1308", "RAG Chunks", BLUE)
    add_metric(slide, 10.75, 6.0, "1024", "BGE-M3向量维度", GREEN)
    add_metric(slide, 15.5, 6.0, "44", "自动化测试通过", AMBER)
    add_textbox(slide, "MinerU · FAISS · BGE-M3 · BGE-Reranker · Tavily · Vue/FastAPI", 1.3, 9.1, 20, 0.6, size=13, color=MUTED)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "01", "本阶段目标与完成概览", "围绕“可溯源知识库 + RAG检索 + 非专业入口”完善端到端生成链路")
    add_card(slide, 1.2, 3.0, 6.9, 3.2, "知识底座", "使用 MinerU 解析 docs 中焊接标准 PDF，生成面向向量库的 JSONL chunk。表格 chunk 保留文件名、页码、表头、上下文、标题和脚注。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 3.2, "检索增强", "建立 FAISS 向量库；RAG 检索同时包含 BGE-M3 语义召回和 BM25 关键词召回，RRF 融合后使用 BGE-Reranker 精排。", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 3.2, "产品入口", "前端新增非专业描述模式，用户可输入自然语言需求，由新智能体自动识别必填字段并进入标准文档生成流程。", GREEN)
    add_card(slide, 1.2, 6.8, 10.4, 2.55, "当前可运行链路", "非专业描述 / 专业问答 → 必填字段校验 → RAG优先标准制定 → DOCX模板填充 → 下载生成文档", AMBER)
    add_card(slide, 12.3, 6.8, 10.4, 2.55, "质量保障", "后端回归测试 44 项通过；前端 npm run build 通过；模型服务异常时非专业解析可退化为规则抽取，避免前端白屏。", TEAL)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "02", "PDF知识库构建：MinerU解析与Chunk设计", "从分散标准文档到可向量检索、可追溯的结构化知识资产")
    add_flow(slide, [("PDF输入", "docs/*.pdf\n11份标准/规范"), ("MinerU解析", "content_list_v2\nMarkdown/中间JSON"), ("Chunk生成", "文本/表格分块\n上下文窗口"), ("溯源字段", "文件名/页码/表头\ncaption/footnote"), ("JSONL输出", "data/rag/pdf_chunks.jsonl\nby_file旁路")], 3.0)
    add_card(slide, 1.2, 6.4, 6.9, 2.55, "表格重点处理", "大表可按行拆分，每个子表重复携带完整表头和相同溯源信息，保证后续 RAG 命中表格片段时能回到原文位置。", TEAL)
    add_card(slide, 8.5, 6.4, 6.9, 2.55, "上下文增强", "表格 chunk 中加入表格前后 1-3 个文本块，减少孤立表格无法判断适用条件的问题。", BLUE)
    add_card(slide, 15.8, 6.4, 6.9, 2.55, "产物管理", "保留 MinerU 原始输出、按文件 chunks、summary 与 manifest，便于审计、重跑和问题定位。", GREEN)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "03", "向量入库与混合检索", "FAISS 不作为独立服务启动，运行时由检索程序直接读取本地索引文件")
    add_metric(slide, 1.2, 3.0, "1308", "向量条目", TEAL)
    add_metric(slide, 6.0, 3.0, "1024", "向量维度", BLUE)
    add_metric(slide, 10.8, 3.0, "Top10", "BGE-M3粗召回", GREEN)
    add_metric(slide, 15.6, 3.0, "Top5", "Reranker精排", AMBER)
    add_flow(slide, [("Query", "用户字段/问题"), ("BGE-M3", "语义召回 Top10"), ("BM25", "关键词召回 Top10"), ("RRF", "k=60 融合"), ("Reranker", "精排 Top5")], 6.15)
    add_textbox(slide, "索引文件：data/rag/faiss/docs.faiss；元数据：data/rag/faiss/chunks.jsonl。", 1.3, 9.15, 20, 0.55, size=12.5, color=MUTED)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "04", "RAG前置到标准制定智能体", "把本地标准库作为第一证据源，网络搜索和YAML转为补充/兜底")
    add_flow(slide, [("输入字段", "工艺/对象/接头\n母材/尺寸"), ("本地RAG", "标准PDF证据\n表格上下文"), ("Tavily", "网络补充\n相似案例"), ("YAML", "行业基准\n规则兜底"), ("字段输出", "document_fields")], 3.0)
    add_card(slide, 1.2, 6.45, 6.9, 2.65, "字段优先级", "用户硬约束 > RAG证据 > 可信网络证据 > 行业YAML > 规则默认值 > /", TEAL)
    add_card(slide, 8.5, 6.45, 6.9, 2.65, "表格填充策略", "填充焊材、预热、层间温度、电流、电压、焊接速度等字段时优先依据 RAG 命中的标准表格和上下文。", BLUE)
    add_card(slide, 15.8, 6.45, 6.9, 2.65, "LLM约束", "Prompt 已同步要求只能在 RAG、网络、YAML、规则候选中选择或归纳短值，避免无证据参数。", GREEN)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "05", "新增智能体：非专业描述解析", "让用户不懂专业字段也能从自然语言需求进入标准文档生成")
    add_card(slide, 1.2, 3.0, 6.9, 3.2, "输入方式", "用户直接输入描述性文字，例如：我有两块20mm厚的20号钢板，需要帮我生成一份焊接工艺。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 3.2, "解析能力", "规则抽取明确事实；结合 RAG/Tavily 证据和可选 LLM 归纳补齐 welding_reask_agent 的 5 个必填字段。", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 3.2, "校验策略", "非专业模式只对必填字段是否为空做校验；识别后用户仍可在右侧字段面板进行人工修正。", GREEN)
    add_card(slide, 1.2, 6.8, 10.4, 2.45, "新增API", "POST /api/nonprofessional/analyze：返回识别字段与缺失问题\nPOST /api/nonprofessional/generate：完整字段直接进入标准文档生成", AMBER)
    add_card(slide, 12.3, 6.8, 10.4, 2.45, "端到端链路", "NaturalLanguageWeldingAgent → welding_reask_agent → welding_standard_agent → welding_document_agent", TEAL)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "06", "焊接工艺同义词标准化", "自然语言可以多样，后续标准制定必须使用统一工艺代号")
    add_card(slide, 1.2, 3.0, 6.9, 2.25, "SMAW", "焊条电弧焊 / 手工电弧焊 / 手弧焊 / manual metal arc / MMA", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 2.25, "GTAW", "氩弧焊 / 钨极氩弧焊 / 钨极气体保护焊 / TIG", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 2.25, "GMAW", "气保焊 / 二保焊 / CO2气体保护焊 / MIG / MAG", GREEN)
    add_card(slide, 1.2, 5.8, 6.9, 2.25, "FCAW", "药芯焊丝电弧焊 / 药芯焊丝气保焊", AMBER)
    add_card(slide, 8.5, 5.8, 6.9, 2.25, "SAW", "埋弧焊", TEAL)
    add_card(slide, 15.8, 5.8, 6.9, 2.25, "GTAW+SMAW", "氩弧焊打底、焊条电弧焊填充等组合描述自动归一为复合工艺", BLUE)
    add_textbox(slide, "前端焊接工艺字段固定为标准下拉选项，后端 API 返回前已完成归一化。", 1.3, 8.95, 20.5, 0.55, size=12.5, color=MUTED)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "07", "前端模块升级", "保留专业问答模式，新增非专业描述模式，一套字段面板支撑两种入口")
    add_card(slide, 1.2, 3.0, 6.9, 3.15, "模式切换", "顶部增加分段控件：专业问答模式 / 非专业描述模式。现有多轮追问流程不受影响。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 3.15, "非专业工作区", "左侧大文本框输入需求描述，提供“分析描述”和“生成标准文档”两个核心动作。", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 3.15, "字段面板", "右侧展示自动识别结果。焊接工艺固定标准下拉，其余非专业字段可自由编辑。", GREEN)
    add_bullets(slide, ["生成时复用现有进度条与下载链接区域", "字段不完整时返回明确缺失项，不生成空文档", "前端构建已通过 vue-tsc 与 Vite production build"], 1.4, 6.8, 20.5, 2.0, size=13)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "08", "当前系统架构", "从输入、知识检索、字段生成到文档落地的多智能体协同")
    add_flow(slide, [("用户入口", "专业问答\n非专业描述"), ("字段层", "重问校验\n自然语言解析"), ("知识层", "MinerU + FAISS\nTavily + YAML"), ("标准层", "RAG优先字段生成\nLLM受约束归纳"), ("交付层", "DOCX模板填充\n下载文档")], 3.0)
    add_card(slide, 1.2, 6.45, 21.5, 2.7, "工程模块", "src/pipeline_welding/agents：四类智能体；src/pipeline_welding/rag：混合检索；src/pipeline_welding/api：FastAPI接口；frontend/src：Vue工作台；configs：模型、RAG、MCP与行业规则配置。", TEAL)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "09", "端到端示例", "从一句非专业描述到可下载的焊接工艺文档")
    add_card(slide, 1.2, 3.0, 6.9, 2.7, "用户描述", "我有两块20mm厚的20号钢板，想用焊条电弧焊生成焊接工艺。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 2.7, "自动字段", "welding_process: SMAW\nwelding_object: 板材\njoint_type: 对接\nbase_material: 20#\nbase_thickness_or_diameter: 板厚 20 mm", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 2.7, "后续生成", "标准制定智能体优先使用本地 RAG 证据，再由网络/YAML补齐，最终写入 MHPWPS-062.docx。", GREEN)
    add_flow(slide, [("分析描述", "字段识别"), ("人工修正", "可选"), ("标准制定", "RAG优先"), ("模板填充", "DOCX"), ("下载", "结果文件")], 6.9)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "10", "测试与验证", "用自动化回归保证新链路不破坏旧功能")
    add_metric(slide, 1.2, 3.0, "44", "pytest回归测试", TEAL)
    add_metric(slide, 6.0, 3.0, "通过", "python compileall", BLUE)
    add_metric(slide, 10.8, 3.0, "通过", "npm run build", GREEN)
    add_metric(slide, 15.6, 3.0, "通过", "RAG脚本CLI", AMBER)
    add_card(slide, 1.2, 6.05, 6.9, 2.7, "后端覆盖", "自然语言解析、工艺同义词归一、API analyze/generate、RAG候选优先级、文档字段填充。", TEAL)
    add_card(slide, 8.5, 6.05, 6.9, 2.7, "前端验证", "Vue 类型检查和生产构建通过；专业模式与非专业模式共享字段面板但行为分离。", BLUE)
    add_card(slide, 15.8, 6.05, 6.9, 2.7, "真实检索冒烟", "BGE-M3 与 BGE-Reranker 服务在线时，示例查询可命中 GB/T 31032 表格片段并返回页码、表头和上下文。", GREEN)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "11", "应用价值与阶段收益", "当前系统已经从“规则问答”升级为“知识增强的文档生成工作台”")
    add_card(slide, 1.2, 3.0, 6.9, 3.15, "降低使用门槛", "非专业用户无需理解全部字段，描述需求即可得到可编辑字段，减少沟通轮次。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 3.15, "增强可追溯性", "标准字段来自本地 PDF 知识库、网络证据和 YAML 规则，表格 chunk 保留页码与表头。", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 3.15, "提升生成质量", "RAG优先、RRF融合、Reranker精排和 Prompt 约束共同降低无证据参数和模板污染风险。", GREEN)
    add_card(slide, 1.2, 6.8, 21.5, 2.15, "阶段结论", "系统已具备离线知识库构建、混合检索、字段智能补齐、标准化工艺归一、DOCX自动生成和前端双模式交互能力，可进入更多样例验证与业务规则精修阶段。", AMBER)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "12", "下一步计划", "围绕数据质量、工程可用性和业务可信度继续增强")
    add_card(slide, 1.2, 3.0, 6.9, 3.3, "数据层", "补充企业历史 WPS/PQR 样本；对 MinerU 解析失败或图片型表格进行人工复核；沉淀字段级来源标注。", TEAL)
    add_card(slide, 8.5, 3.0, 6.9, 3.3, "检索层", "加入查询改写和字段定向检索；根据表格/文本类型调权；优化大表 chunk 的召回粒度。", BLUE)
    add_card(slide, 15.8, 3.0, 6.9, 3.3, "产品层", "增加生成结果预览、字段来源展示、人工确认记录和导出前检查清单。", GREEN)
    add_card(slide, 1.2, 6.9, 21.5, 2.1, "近期重点", "用典型材料和工艺场景做批量案例测试：20#钢板、Q345管道、A106 Gr.B、L245/L360管线钢，验证字段准确率和人工修正成本。", AMBER)

    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_textbox(slide, "谢谢", 1.2, 2.2, 8, 1.0, size=40, bold=True, color=NAVY)
    add_textbox(slide, "PipelineWelding · 管道焊接标准生成多智能体系统", 1.3, 4.0, 14, 0.6, size=16, color=TEAL)
    add_textbox(slide, "当前进展：MinerU知识库 · FAISS混合检索 · RAG优先标准制定 · 非专业描述入口", 1.3, 5.0, 18, 0.6, size=13, color=TEXT)
    add_textbox(slide, "厚德 · 博学 · 求是", 1.3, 8.8, 8, 0.55, size=14, color=MUTED)

    prs.save(str(OUTPUT))


if __name__ == "__main__":
    build_ppt()
    print(OUTPUT)
